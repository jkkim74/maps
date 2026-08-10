"""Build the annotated MAPS stock-analysis screen-design PowerPoint deck."""

from __future__ import annotations

import argparse
import subprocess
import tempfile
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_PROTOTYPE = ROOT / "docs" / "ui-design" / "maps-analysis-trade-prototype.html"
DEFAULT_ASSETS = ROOT / "docs" / "ui-design" / "assets"
DEFAULT_OUTPUT = ROOT / "docs" / "ui-design" / "MAPS_종목분석_전략매매_화면설계서.pptx"
CHROME_CANDIDATES = (
    Path(r"C:\Program Files\Google\Chrome\Application\chrome.exe"),
    Path(r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"),
    Path(r"C:\Program Files\Microsoft\Edge\Application\msedge.exe"),
)

BG = RGBColor(7, 13, 23)
PANEL = RGBColor(16, 29, 48)
WHITE = RGBColor(235, 242, 251)
MUTED = RGBColor(151, 166, 187)
BLUE = RGBColor(59, 130, 246)
GREEN = RGBColor(34, 197, 94)
ORANGE = RGBColor(245, 158, 11)
RED = RGBColor(239, 68, 68)


@dataclass(frozen=True)
class SlideSpec:
    """One slide's screen capture and user/system responsibility mapping."""

    number: int
    title: str
    subtitle: str
    screen: str | None
    scenario: str
    user_actions: tuple[str, ...]
    system_actions: tuple[str, ...]
    result: str


SLIDES: tuple[SlideSpec, ...] = (
    SlideSpec(1, "종목 분석에서 자동매수까지", "MAPS 종목 분석–전략매매 화면설계서", None, "normal", ("화면설계 목적과 범위를 확인한다.",), ("분석·설정·감시·주문의 연결 구조를 제공한다.",), "실제 주문이 없는 화면설계용 산출물입니다."),
    SlideSpec(2, "화면설계서 읽는 법", "파란색 U는 사용자, 초록색 S는 시스템", None, "normal", ("U 번호의 선택·입력·확인 작업을 수행한다.",), ("S 번호의 조회·계산·검증·저장을 자동 수행한다.",), "각 화면에서 누가 무엇을 하는지 같은 번호로 연결합니다."),
    SlideSpec(3, "전체 사용자 여정", "분석부터 자동청산까지의 역할 분리", "journey", "normal", ("종목을 분석하고 매매 방식을 선택한다.", "가격·금액·위험을 확인해 실행한다."), ("AI 전략과 안전한도를 계산한다.", "ARMED 저장 후 현재가와 주문을 감시한다."), "분석 → 설정 → 무장 → 순차매수 → 전체청산"),
    SlideSpec(4, "AS-IS와 TO-BE", "수동 복사 흐름을 구조화 전략으로 연결", "asis", "normal", ("기존에는 AI 문장을 보고 값을 다시 입력했다.", "새 화면에서는 제안값을 검토하고 승인한다."), ("AI 숫자전략을 검증해 매매 설정에 자동 입력한다.",), "분석 근거와 워치 전략이 하나의 흐름으로 이어집니다."),
    SlideSpec(5, "종목 검색과 분석 실행", "사용자가 종목을 정하고 시스템이 분석", "analysis-input", "normal", ("종목명 또는 코드를 입력한다.", "분석 실행을 누른다."), ("종목을 확인하고 기술·재무·AI 분석을 수행한다.",), "분석 결과 화면으로 이동합니다."),
    SlideSpec(6, "분석 결과와 AI 전략", "설명과 주문용 숫자를 분리", "analysis-result", "normal", ("분석과 BUY 의견을 검토한다.", "매매 설정을 누른다."), ("3개 진입가와 목표·손절을 구조화해 검증한다.",), "BUY이면 자동값을 제공하고 매매 설정으로 이동합니다."),
    SlideSpec(7, "AI 실패·관망 시 수동 전환", "AI 의견을 자동주문으로 강제하지 않음", "ai-fallback", "ai-fail", ("경고를 확인하고 수동입력 여부를 결정한다.",), ("자동값을 비우고 수동값에도 같은 검증을 적용한다.",), "AI 실패나 WATCH가 자동매수로 연결되지 않습니다."),
    SlideSpec(8, "매매 방식 선택", "단일 또는 3분할을 명시 선택", "trade-mode", "normal", ("단일매매 또는 3분할매매를 선택한다.",), ("선택 전 실행을 막고 모드별 입력폼을 준비한다.",), "단일은 1회, 분할은 30/30/40의 3회 계획입니다."),
    SlideSpec(9, "단일매매 설정", "한 번의 진입 가격과 수량", "single-setup", "normal", ("금액·진입·목표·손절을 확인한다.", "필요하면 총액을 감액한다."), ("안전한도와 정수주 수량을 계산한다.", "가격과 손익비를 검증한다."), "검증된 단일 회차가 최종 확인으로 전달됩니다."),
    SlideSpec(10, "3분할매매 설정", "회차별 가격과 30/30/40 배분", "split-setup", "normal", ("총액과 3개 진입가·비중을 확인한다.", "목표·손절을 검토하고 최종 확인을 누른다."), ("브로커 현금과 안전한도를 조회한다.", "회차별 정수주 수량과 주문금액을 계산한다."), "한 주기에 한 회차만 주문하는 계획이 만들어집니다."),
    SlideSpec(11, "안전 최대금액 계산", "주문가능 현금과 투자 허용액의 구분", "safe-budget", "normal", ("제안금액을 확인하고 감액할 수 있다.",), ("현금·종목노출·포트폴리오·손절위험 중 최솟값을 제시한다.",), "안전 최대금액보다 증액할 수 없습니다."),
    SlideSpec(12, "입력 검증과 실행 차단", "오류가 남아 있으면 주문 준비 불가", "validation", "gate-off", ("오류 메시지에 따라 값이나 운영설정을 수정한다.",), ("가격순서·비중·호가·잔고·중복·게이트를 검증한다.",), "모든 오류가 해소될 때만 최종 확인이 활성화됩니다."),
    SlideSpec(13, "REAL/PAPER 최종 확인", "계좌·금액·손절을 실행 직전 재확인", "confirm", "normal", ("계좌 모드와 최대금액을 확인한다.", "매매 실행·즉시 무장을 누른다."), ("잔고와 안전 게이트를 다시 검증한다.", "계획 저장과 ARMED 전환을 함께 처리한다."), "주문을 즉시 보내지 않고 워치 감시를 시작합니다."),
    SlideSpec(14, "즉시 무장 완료", "첫 자동매수 조건을 안내", "armed", "normal", ("다음 진입 조건과 만료일을 확인한다.", "워치리스트로 이동한다."), ("60초마다 현재가를 조회해 조건 충족 시 지정가 주문한다.",), "ARMED 상태로 자동매수 감시가 시작됩니다."),
    SlideSpec(15, "분석 워치리스트", "진행률과 다음 조건을 한눈에 확인", "watchlist", "normal", ("체결 회차·다음 가격·보류 사유를 확인한다.", "종목을 선택해 상세를 연다."), ("현재가·평균체결가·주문상태를 동기화한다.",), "단일/분할과 0/3~3/3 진행상태를 표시합니다."),
    SlideSpec(16, "분할 회차 상세", "부분체결과 남은 매수 관리", "leg-detail", "partial", ("회차별 체결과 잔량을 확인한다.", "필요하면 남은 매수를 중단한다."), ("직전 회차 체결 후 다음 회차만 허용한다.", "부분체결 잔량만 재주문하고 보유분 청산은 계속 감시한다."), "추가매수와 기존 보유분 청산 관리가 분리됩니다."),
    SlideSpec(17, "모바일 워치리스트", "웹 전략의 진행상태 확인", "mobile", "normal", ("진행률과 다음 가격을 확인한다.", "필요하면 남은 매수를 중단한다."), ("웹과 같은 회차·체결·보류 정보를 제공한다.",), "모바일은 조회·중단, 새 전략 설정은 웹에서 수행합니다."),
    SlideSpec(18, "상태 흐름과 검수 기준", "ARMED → BOUGHT → CLOSED", "state-flow", "normal", ("각 상태의 의미와 가능한 행동을 확인한다.",), ("만료는 신규매수만 막고 BOUGHT 청산은 계속 감시한다.",), "신규 진입 차단이 보유분 청산 중단으로 이어지지 않습니다."),
)


def find_browser(explicit: Path | None = None) -> Path:
    """Return a supported Chromium browser executable."""
    if explicit and explicit.is_file():
        return explicit
    for candidate in CHROME_CANDIDATES:
        if candidate.is_file():
            return candidate
    raise FileNotFoundError("Chrome or Edge executable was not found")


def capture_screens(prototype: Path, assets_dir: Path, browser: Path) -> None:
    """Capture every slide-backed HTML screen at a deterministic size."""
    if not prototype.is_file():
        raise FileNotFoundError(prototype)
    assets_dir.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="maps-ppt-browser-") as profile:
        for spec in SLIDES:
            if spec.screen is None:
                continue
            output = (assets_dir / f"{spec.number:02d}.png").resolve()
            url = f"{prototype.resolve().as_uri()}?screen={spec.screen}&scenario={spec.scenario}&capture=1"
            command = [
                str(browser),
                "--headless=new",
                "--disable-gpu",
                "--hide-scrollbars",
                "--no-first-run",
                f"--user-data-dir={profile}",
                "--window-size=1120,740",
                "--force-device-scale-factor=1",
                f"--screenshot={output}",
                url,
            ]
            completed = subprocess.run(command, check=False, capture_output=True, text=True, timeout=30)
            if completed.returncode != 0 or not output.is_file() or output.stat().st_size == 0:
                raise RuntimeError(
                    f"Browser capture failed for slide {spec.number}: "
                    f"{completed.stderr or completed.stdout}"
                )


def _set_background(slide, color: RGBColor = BG) -> None:
    """Set a slide's solid background color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: int = 16,
    color: RGBColor = WHITE,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    """Add a consistently styled text box using inch coordinates."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(0.04)
    frame.margin_top = frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = "Malgun Gothic"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color


def _add_panel(slide, left: float, top: float, width: float, height: float, color: RGBColor) -> None:
    """Add a flat rounded-looking content panel."""
    shape = slide.shapes.add_shape(5, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(43, 62, 87)


def _add_action_panel(
    slide,
    heading: str,
    marker: str,
    actions: tuple[str, ...],
    left: float,
    top: float,
    color: RGBColor,
) -> None:
    """Add one user/system responsibility panel with numbered bullets."""
    height = 1.52 if len(actions) <= 2 else 1.82
    _add_panel(slide, left, top, 4.12, height, PANEL)
    _add_text(slide, marker, left + 0.12, top + 0.12, 0.42, 0.38, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    circle = slide.shapes[-1]
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.color.rgb = color
    _add_text(slide, heading, left + 0.62, top + 0.1, 3.25, 0.35, size=12, bold=True)
    bullet_text = "\n".join(f"{index}. {action}" for index, action in enumerate(actions, 1))
    _add_text(slide, bullet_text, left + 0.18, top + 0.5, 3.72, height - 0.58, size=9, color=MUTED)


def _add_slide_header(slide, spec: SlideSpec) -> None:
    """Add number, title, subtitle, and sample-data label."""
    _add_text(slide, f"{spec.number:02d}", 0.28, 0.18, 0.55, 0.42, size=15, color=BLUE, bold=True)
    _add_text(slide, spec.title, 0.86, 0.12, 7.5, 0.48, size=22, bold=True)
    _add_text(slide, spec.subtitle, 0.88, 0.53, 7.4, 0.27, size=9, color=MUTED)
    _add_text(slide, "예시 데이터 · 실제 주문 없음", 10.6, 0.2, 2.45, 0.3, size=9, color=ORANGE, bold=True, align=PP_ALIGN.RIGHT)


def _build_cover(slide, spec: SlideSpec) -> None:
    """Build the title slide."""
    _set_background(slide)
    _add_text(slide, "MAPS", 0.7, 0.62, 2.0, 0.5, size=17, color=BLUE, bold=True)
    _add_text(slide, spec.title, 0.7, 2.0, 11.8, 1.0, size=34, bold=True)
    _add_text(slide, spec.subtitle, 0.74, 3.05, 9.5, 0.5, size=18, color=MUTED)
    _add_panel(slide, 0.74, 4.12, 11.85, 1.55, PANEL)
    _add_text(slide, "HTML 프로토타입", 1.05, 4.45, 2.7, 0.42, size=15, color=BLUE, bold=True)
    _add_text(slide, "→", 3.8, 4.43, 0.55, 0.42, size=18, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "U/S 역할 주석", 4.45, 4.45, 2.55, 0.42, size=15, color=GREEN, bold=True)
    _add_text(slide, "→", 7.1, 4.43, 0.55, 0.42, size=18, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "18장 화면설계서", 7.8, 4.45, 3.2, 0.42, size=15, color=WHITE, bold=True)
    _add_text(slide, "화면설계용 예시 데이터 · 실제 MAPS API 및 주문 호출 없음", 0.74, 6.75, 8.0, 0.3, size=9, color=ORANGE)


def _build_guide(slide, spec: SlideSpec) -> None:
    """Build the user/system legend slide."""
    _set_background(slide)
    _add_slide_header(slide, spec)
    _add_panel(slide, 0.7, 1.2, 5.8, 4.75, PANEL)
    _add_text(slide, "U", 1.05, 1.62, 0.7, 0.7, size=22, bold=True, align=PP_ALIGN.CENTER)
    badge = slide.shapes[-1]
    badge.fill.solid(); badge.fill.fore_color.rgb = BLUE; badge.line.color.rgb = BLUE
    _add_text(slide, "사용자가 하는 일", 1.95, 1.62, 3.4, 0.7, size=21, bold=True)
    _add_text(slide, "직접 선택 · 입력 · 검토 · 최종 승인\n\n화면의 파란색 U 번호와 오른쪽 설명이 연결됩니다.", 1.05, 2.65, 4.8, 1.9, size=14, color=MUTED)
    _add_panel(slide, 6.82, 1.2, 5.8, 4.75, PANEL)
    _add_text(slide, "S", 7.18, 1.62, 0.7, 0.7, size=22, bold=True, align=PP_ALIGN.CENTER)
    badge = slide.shapes[-1]
    badge.fill.solid(); badge.fill.fore_color.rgb = GREEN; badge.line.color.rgb = GREEN
    _add_text(slide, "시스템이 해주는 일", 8.08, 1.62, 3.65, 0.7, size=21, bold=True)
    _add_text(slide, "자동 조회 · 계산 · 검증 · 저장 · 감시\n\n화면의 초록색 S 번호와 오른쪽 설명이 연결됩니다.", 7.18, 2.65, 4.8, 1.9, size=14, color=MUTED)
    _add_text(slide, spec.result, 0.85, 6.35, 11.7, 0.48, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def _build_screen_slide(slide, spec: SlideSpec, image_path: Path) -> None:
    """Build one screen slide with an HTML capture and native annotations."""
    _set_background(slide)
    _add_slide_header(slide, spec)
    if not image_path.is_file():
        raise FileNotFoundError(image_path)
    slide.shapes.add_picture(str(image_path), Inches(0.28), Inches(0.88), width=Inches(8.62), height=Inches(5.68))
    _add_action_panel(slide, "사용자가 하는 일", "U", spec.user_actions, 9.05, 1.05, BLUE)
    system_top = 2.8 if len(spec.user_actions) <= 2 else 3.08
    _add_action_panel(slide, "시스템이 해주는 일", "S", spec.system_actions, 9.05, system_top, GREEN)
    _add_panel(slide, 9.05, 5.48, 4.12, 1.08, RGBColor(12, 23, 39))
    _add_text(slide, "처리 결과", 9.22, 5.58, 1.2, 0.28, size=10, color=ORANGE, bold=True)
    _add_text(slide, spec.result, 9.22, 5.87, 3.72, 0.52, size=9, color=WHITE)
    _add_text(slide, "입력 → 시스템 처리 → 결과 상태", 0.38, 6.75, 3.0, 0.3, size=9, color=MUTED)


def build_presentation(output: Path, captures_dir: Path) -> None:
    """Create the final 18-slide 16:9 PowerPoint presentation."""
    deck = Presentation()
    deck.slide_width = Inches(13.333333)
    deck.slide_height = Inches(7.5)
    blank = deck.slide_layouts[6]
    for spec in SLIDES:
        slide = deck.slides.add_slide(blank)
        if spec.number == 1:
            _build_cover(slide, spec)
        elif spec.number == 2:
            _build_guide(slide, spec)
        else:
            _build_screen_slide(slide, spec, captures_dir / f"{spec.number:02d}.png")
    output.parent.mkdir(parents=True, exist_ok=True)
    deck.save(output)


def parse_args() -> argparse.Namespace:
    """Parse command-line options for reproducible document generation."""
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--prototype", type=Path, default=DEFAULT_PROTOTYPE)
    parser.add_argument("--assets", type=Path, default=DEFAULT_ASSETS)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument("--browser", type=Path)
    parser.add_argument("--skip-capture", action="store_true")
    return parser.parse_args()


def main() -> int:
    """Capture prototype screens and build the PowerPoint deck."""
    args = parse_args()
    if not args.skip_capture:
        capture_screens(args.prototype, args.assets, find_browser(args.browser))
    build_presentation(args.output, args.assets)
    print(f"Created {args.output} ({len(SLIDES)} slides)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
