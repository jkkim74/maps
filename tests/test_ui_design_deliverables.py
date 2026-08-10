"""Contract tests for the standalone UI prototype and its PPT deliverable."""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
PROTOTYPE = ROOT / "docs" / "ui-design" / "maps-analysis-trade-prototype.html"
CAPTURES = ROOT / "docs" / "ui-design" / "assets"
PRESENTATION = ROOT / "docs" / "ui-design" / "MAPS_종목분석_전략매매_화면설계서.pptx"


def test_prototype_contains_all_annotated_screens_without_live_network_calls() -> None:
    """The prototype must be complete, annotated, and incapable of live trading."""
    html = PROTOTYPE.read_text(encoding="utf-8")
    required_screens = {
        "journey",
        "asis",
        "analysis-input",
        "analysis-result",
        "ai-fallback",
        "trade-mode",
        "single-setup",
        "split-setup",
        "safe-budget",
        "validation",
        "confirm",
        "armed",
        "watchlist",
        "leg-detail",
        "mobile",
        "state-flow",
    }

    for screen in required_screens:
        assert f'"{screen}"' in html
    assert "사용자가 하는 일" in html
    assert "시스템이 해주는 일" in html
    assert "예시 데이터 · 실제 주문이 발생하지 않습니다" in html
    assert "marker user" in html
    assert "marker system" in html
    assert "30 / 30 / 40" in html
    assert "fetch(" not in html
    assert "XMLHttpRequest" not in html
    assert "WebSocket" not in html


def test_prototype_blocks_unsafe_amounts_and_preserves_strategy_mode() -> None:
    """The demonstrator must not present an unsafe or mode-changing happy path."""
    html = PROTOTYPE.read_text(encoding="utf-8")

    assert 'const brokerCash=state.scenario==="cash"?2500000:12500000' in html
    assert 'data-role="confirm-trade" ${over?"disabled":""}' in html
    assert "const tradeBlocks=()=>" in html
    assert 'data-role="arm-strategy" ${blocked?"disabled":""}' in html
    assert 'state.tradeMode==="single"?"단일":"3분할"' in html
    assert "applyManualStrategy()" in html
    assert "manualDraft" in html
    assert "const splitPlan=()=>" in html
    assert 'params.get("mode")' in html
    assert "9,856,000원" not in html
    assert "<td>62주</td>" not in html


def test_checked_in_captures_are_real_browser_screens() -> None:
    """The generated deck must be backed by the checked-in, non-blank captures."""
    captures = sorted(CAPTURES.glob("*.png"))
    assert [path.name for path in captures] == [f"{number:02d}.png" for number in range(3, 19)]
    for capture in captures:
        with Image.open(capture) as image:
            assert image.size == (1120, 740)
            assert all(low < high for low, high in image.convert("RGB").getextrema())


def test_document_generation_dependencies_are_declared() -> None:
    """A fresh project install must be able to run the builder and its tests."""
    requirements = (ROOT / "requirements.txt").read_text(encoding="utf-8").lower()
    assert "pillow" in requirements
    assert "python-pptx" in requirements


def test_checked_in_presentation_embeds_every_checked_in_capture() -> None:
    """The delivered deck must be synchronized with the delivered browser screens."""
    deck = Presentation(PRESENTATION)
    embedded_images = {
        shape.image.blob
        for slide in deck.slides
        for shape in slide.shapes
        if hasattr(shape, "image")
    }
    capture_images = {path.read_bytes() for path in CAPTURES.glob("*.png")}

    assert len(deck.slides) == 18
    assert len(embedded_images) == 16
    assert embedded_images == capture_images


def test_slide_spec_defines_eighteen_user_system_mapped_slides() -> None:
    """Every slide must state the user's action and the system's responsibility."""
    from scripts.build_stock_analysis_ui_ppt import SLIDES

    assert len(SLIDES) == 18
    assert [slide.number for slide in SLIDES] == list(range(1, 19))
    assert SLIDES[0].title == "종목 분석에서 자동매수까지"
    assert SLIDES[-1].title == "상태 흐름과 검수 기준"
    for slide in SLIDES[1:]:
        assert slide.user_actions
        assert slide.system_actions
        assert slide.result


def test_generated_presentation_is_widescreen_and_has_eighteen_slides(tmp_path: Path) -> None:
    """The builder must create a readable 16:9 deck from deterministic captures."""
    from scripts.build_stock_analysis_ui_ppt import SLIDES, build_presentation

    captures = tmp_path / "captures"
    captures.mkdir()
    for slide in SLIDES:
        if slide.screen:
            Image.new("RGB", (1100, 720), "#0b1423").save(captures / f"{slide.number:02d}.png")

    output = tmp_path / "screen-design.pptx"
    build_presentation(output, captures)

    deck = Presentation(output)
    assert len(deck.slides) == 18
    assert round(deck.slide_width / deck.slide_height, 3) == round(16 / 9, 3)
    all_text = "\n".join(shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text"))
    assert "사용자가 하는 일" in all_text
    assert "시스템이 해주는 일" in all_text
    assert "예시 데이터" in all_text
